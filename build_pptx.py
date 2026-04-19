"""
build_pptx.py — Builds a branded 16-slide .pptx from ppt_content.json
+ screenshots fetched from the website-decoder repo via raw GitHub URLs.
"""

import argparse
import json
import sys
from pathlib import Path

import requests
from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE  = RGBColor(0xF4, 0xF5, 0xF7)
DARK_NAVY  = RGBColor(0x16, 0x09, 0x23)
MID_GRAY   = RGBColor(0x33, 0x33, 0x44)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xDD)


def hex_to_rgb(h):
    try:
        h = h.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except Exception:
        return DARK_NAVY


def is_dark(h):
    try:
        h = h.lstrip("#")
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        return (0.299 * r + 0.587 * g + 0.114 * b) / 255 < 0.45
    except Exception:
        return True


def lighten(h, factor=0.4):
    try:
        h = h.lstrip("#")
        r = int(int(h[0:2], 16) * (1 - factor) + 255 * factor)
        g = int(int(h[2:4], 16) * (1 - factor) + 255 * factor)
        b = int(int(h[4:6], 16) * (1 - factor) + 255 * factor)
        return RGBColor(r, g, b)
    except Exception:
        return LIGHT_GRAY


def fetch_bytes(url, retries=3):
    for attempt in range(retries):
        try:
            r = requests.get(url, timeout=25)
            r.raise_for_status()
            return r.content
        except Exception as e:
            if attempt < retries - 1:
                import time; time.sleep(10)
            else:
                print(f"  ! Failed: {url}: {e}")
    return None


def fetch_screenshot(base_url, filename, cache_dir):
    # If cache_dir IS the screenshots dir (local mode), just check directly
    candidate = Path(cache_dir) / filename
    if candidate.exists():
        return str(candidate)
    # Remote mode: download and cache
    cache_dir = Path(cache_dir)
    cache_dir.mkdir(parents=True, exist_ok=True)
    data = fetch_bytes(f"{base_url}/{filename}")
    if data is None:
        return None
    candidate.write_bytes(data)
    return str(candidate)


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, hex_color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = hex_to_rgb(hex_color)


def add_rect(slide, left, top, w, h, fill_hex, line_hex=None):
    s = slide.shapes.add_shape(1, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = hex_to_rgb(fill_hex)
    if line_hex:
        s.line.color.rgb = hex_to_rgb(line_hex)
        s.line.width = Pt(0.5)
    else:
        s.line.fill.background()
    return s


def add_text(slide, text, left, top, w, h,
             size=Pt(14), bold=False, color=None,
             align=PP_ALIGN.LEFT, wrap=True):
    if color is None:
        color = WHITE
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    return tf


def add_text_with_links(slide, text, left, top, w, h,
                        size=Pt(13), color=None, post_urls=None):
    """Add text where 'Post N' references are hyperlinked to LinkedIn URLs."""
    import re
    if color is None:
        color = MID_GRAY
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]

    parts = re.split(r'(Post \d+)', text)
    for part in parts:
        m = re.match(r'Post (\d+)', part)
        run = p.add_run()
        run.text = part
        run.font.size = size
        run.font.color.rgb = color if isinstance(color, RGBColor) else color
        if m and post_urls:
            num = int(m.group(1))
            url = post_urls.get(num)
            if url:
                run.hyperlink.address = url
                run.font.underline = True
                run.font.color.rgb = RGBColor(0x53, 0x9B, 0xD9)
    return tf


def add_image(slide, path, left, top, w, h, accent_hex="#555566"):
    if path and Path(path).exists():
        try:
            with PILImage.open(path) as im:
                im.verify()
            slide.shapes.add_picture(path, left, top, width=w, height=h)
            return True
        except Exception as e:
            print(f"  ! Image error: {e}")
    s = add_rect(slide, left, top, w, h, accent_hex)
    tf = s.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = "[ screenshot ]"
    r.font.color.rgb = WHITE
    r.font.size = Pt(11)
    return False


def slide_footer(slide, n, domain, accent_hex, on_dark=False):
    gray = lighten(accent_hex, 0.5) if on_dark else RGBColor(0xAA, 0xAA, 0xBB)
    add_text(slide, domain.upper(),
             Inches(0.5), SLIDE_H - Inches(0.38), Inches(4), Inches(0.35),
             size=Pt(8), bold=True, color=gray)
    add_text(slide, str(n),
             SLIDE_W - Inches(0.55), SLIDE_H - Inches(0.38),
             Inches(0.4), Inches(0.35),
             size=Pt(9), color=gray, align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────────────────────────────────────
# LAYOUT BUILDERS
# ─────────────────────────────────────────────────────────────────────────────

def build_cover(prs, s, img_path, n, dark_hex, accent_hex):
    slide = blank_slide(prs)
    fill_bg(slide, dark_hex)

    # Right screenshot
    add_image(slide, img_path,
              Inches(5.0), Inches(0.15),
              Inches(4.85), Inches(5.3), accent_hex)

    # Subtle dark overlay on right side edge for blending
    add_rect(slide, Inches(4.85), Inches(0), Inches(0.3), SLIDE_H, dark_hex)

    # Left accent stripe
    add_rect(slide, Inches(0.32), Inches(0.85), Pt(3.5), Inches(3.2), accent_hex)

    # Title
    add_text(slide, s.get("title", "Website Decode"),
             Inches(0.55), Inches(0.85), Inches(4.1), Inches(1.9),
             size=Pt(36), bold=True, color=WHITE)

    # Domain pill
    domain = s.get("subtitle", "")
    add_rect(slide, Inches(0.55), Inches(2.9), Inches(2.8), Inches(0.48), accent_hex)
    add_text(slide, "  " + domain,
             Inches(0.55), Inches(2.91), Inches(2.8), Inches(0.46),
             size=Pt(15), bold=True, color=WHITE)

    # Deck label
    add_text(slide, "Competitive Intelligence  ·  Website Decode",
             Inches(0.55), Inches(3.55), Inches(4.1), Inches(0.45),
             size=Pt(11), color=lighten(accent_hex, 0.45))

    # Bottom accent bar
    add_rect(slide, Inches(0), SLIDE_H - Pt(5), SLIDE_W, Pt(5), accent_hex)

    slide_footer(slide, n, domain, accent_hex, on_dark=True)


def build_two_column(prs, s, img_path, n, dark_hex, accent_hex):
    """Dark left panel (title + bullets), light right (screenshot)."""
    slide = blank_slide(prs)

    LEFT_W  = Inches(4.15)
    RIGHT_X = Inches(4.3)
    RIGHT_W = SLIDE_W - RIGHT_X

    # Panels
    add_rect(slide, Inches(0), Inches(0), LEFT_W, SLIDE_H, dark_hex)
    add_rect(slide, RIGHT_X, Inches(0), RIGHT_W, SLIDE_H, "#FFFFFF")

    # Top accent bar (full width)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Pt(5), accent_hex)

    # Title
    add_text(slide, s.get("title", ""),
             Inches(0.3), Inches(0.2), Inches(3.65), Inches(0.92),
             size=Pt(20), bold=True, color=WHITE)

    # Accent underline
    add_rect(slide, Inches(0.3), Inches(1.12), Inches(1.4), Pt(2.5), accent_hex)

    # Bullets
    body = s.get("body", [])
    body_color = lighten(dark_hex, 0.78)
    for i, item in enumerate(body):
        top = Inches(1.28) + i * Inches(0.63)
        if top + Inches(0.55) > SLIDE_H - Inches(0.4):
            break
        if item.startswith("✓"):
            dot = "#4CD96E"
            tc  = RGBColor(0x4C, 0xD9, 0x6E)
        elif item.startswith("✗"):
            dot = "#FF6B6B"
            tc  = RGBColor(0xFF, 0x6B, 0x6B)
        else:
            dot = accent_hex
            tc  = body_color
        add_rect(slide, Inches(0.3), top + Inches(0.17), Inches(0.06), Inches(0.22), dot)
        add_text(slide, item,
                 Inches(0.46), top, Inches(3.5), Inches(0.6),
                 size=Pt(12.5), color=tc)

    # Screenshot right
    add_image(slide, img_path,
              RIGHT_X + Inches(0.2), Inches(0.5),
              Inches(5.1), Inches(4.85), accent_hex)

    # Bottom bar
    add_rect(slide, Inches(0), SLIDE_H - Pt(5), LEFT_W, Pt(5), accent_hex)

    slide_footer(slide, n, "", accent_hex, on_dark=False)


def build_bullets(prs, s, n, dark_hex, accent_hex):
    """Dark header band, light body with individual dot bullets."""
    slide = blank_slide(prs)
    HEADER_H = Inches(1.1)

    # Slide background
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, "#F4F5F7")

    # Header band
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, HEADER_H, dark_hex)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Pt(5), accent_hex)

    # Title
    add_text(slide, s.get("title", ""),
             Inches(0.5), Inches(0.15), Inches(9.0), Inches(0.85),
             size=Pt(24), bold=True, color=WHITE)

    # Optional subtitle
    subtitle = s.get("subtitle")
    body_top = HEADER_H + Inches(0.22)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.5), body_top, Inches(9.0), Inches(0.42),
                 size=Pt(14), color=hex_to_rgb(accent_hex))
        body_top += Inches(0.46)

    # Body bullets
    body = s.get("body", [])
    for i, item in enumerate(body):
        top = body_top + i * Inches(0.6)
        if top + Inches(0.5) > SLIDE_H - Inches(0.35):
            break

        if item.startswith("✓"):
            dot = "#4CD96E"
            tc  = RGBColor(0x1A, 0x7A, 0x3C)
        elif item.startswith("✗"):
            dot = "#FF6B6B"
            tc  = RGBColor(0xC0, 0x30, 0x30)
        else:
            dot = accent_hex
            tc  = MID_GRAY

        add_rect(slide, Inches(0.5), top + Inches(0.16), Inches(0.07), Inches(0.26), dot)
        add_text(slide, item,
                 Inches(0.7), top, Inches(9.0), Inches(0.56),
                 size=Pt(13), color=tc)

    add_rect(slide, Inches(0), SLIDE_H - Pt(5), SLIDE_W, Pt(5), accent_hex)
    slide_footer(slide, n, "", accent_hex, on_dark=False)


def build_color_palette(prs, s, n, dark_hex, accent_hex):
    slide = blank_slide(prs)
    fill_bg(slide, dark_hex)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Pt(5), accent_hex)

    add_text(slide, s.get("title", "Brand Color System"),
             Inches(0.5), Inches(0.18), Inches(9.0), Inches(0.78),
             size=Pt(24), bold=True, color=WHITE)
    add_rect(slide, Inches(0.5), Inches(0.95), Inches(1.8), Pt(2.5), accent_hex)

    colors = s.get("colors", [])[:8]
    if not colors:
        return

    swatch_w = min(Inches(1.5), (SLIDE_W - Inches(1.0)) / len(colors) - Inches(0.14))
    swatch_h = Inches(2.25)
    gap      = Inches(0.15)
    total_w  = len(colors) * swatch_w + (len(colors) - 1) * gap
    start_x  = (SLIDE_W - total_w) / 2
    top      = Inches(1.1)

    for i, c in enumerate(colors):
        left    = start_x + i * (swatch_w + gap)
        hex_val = c.get("hex", "#888888")
        role    = c.get("role", "")

        add_rect(slide, left, top, swatch_w, swatch_h, hex_val)

        txt_c = WHITE if is_dark(hex_val) else MID_GRAY
        add_text(slide, hex_val.upper(),
                 left + Inches(0.07), top + swatch_h - Inches(0.42),
                 swatch_w - Inches(0.1), Inches(0.38),
                 size=Pt(10), bold=True, color=txt_c)
        add_text(slide, role,
                 left, top + swatch_h + Inches(0.07),
                 swatch_w, Inches(0.32),
                 size=Pt(10), color=lighten(dark_hex, 0.65),
                 align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(0), SLIDE_H - Pt(5), SLIDE_W, Pt(5), accent_hex)
    slide_footer(slide, n, "", accent_hex, on_dark=True)


def build_quote(prs, s, n, dark_hex, accent_hex):
    slide = blank_slide(prs)
    fill_bg(slide, dark_hex)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Pt(5), accent_hex)

    # Decorative large quote mark
    add_text(slide, "\u201c",
             Inches(0.28), Inches(0.05), Inches(1.8), Inches(1.8),
             size=Pt(96), bold=True, color=lighten(accent_hex, 0.3))

    # Quote
    add_text(slide, s.get("quote", ""),
             Inches(0.85), Inches(0.95), Inches(8.3), Inches(3.3),
             size=Pt(21), bold=False, color=WHITE)

    # Divider + attribution
    add_rect(slide, Inches(0.85), Inches(4.35), Inches(1.5), Pt(2), accent_hex)
    add_text(slide, "— " + s.get("attribution", ""),
             Inches(0.85), Inches(4.5), Inches(8.0), Inches(0.55),
             size=Pt(14), color=hex_to_rgb(accent_hex))

    add_rect(slide, Inches(0), SLIDE_H - Pt(5), SLIDE_W, Pt(5), accent_hex)
    slide_footer(slide, n, s.get("attribution", ""), accent_hex, on_dark=True)


# ─────────────────────────────────────────────────────────────────────────────

SCREENSHOT_LAYOUTS = {"cover", "two-column"}


def build_presentation(content, cache_dir, screenshots_base):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slides = content.get("slides", [])
    cover  = next((s for s in slides if s.get("id") == 1), slides[0] if slides else {})
    dark_hex   = cover.get("bg_color",    "#160923")
    accent_hex = cover.get("accent_color","#7C3AED")

    for i, s in enumerate(slides):
        layout       = s.get("layout", "bullets")
        n            = s.get("id", i + 1)
        slide_accent = s.get("accent_color", accent_hex)

        img_path = None
        if layout in SCREENSHOT_LAYOUTS:
            fname = s.get("screenshot")
            if fname:
                img_path = fetch_screenshot(screenshots_base, fname, cache_dir)

        if layout == "cover":
            build_cover(prs, s, img_path, n, dark_hex, slide_accent)
        elif layout == "two-column":
            build_two_column(prs, s, img_path, n, dark_hex, slide_accent)
        elif layout == "bullets":
            build_bullets(prs, s, n, dark_hex, slide_accent)
        elif layout == "color-palette":
            build_color_palette(prs, s, n, dark_hex, slide_accent)
        elif layout == "quote":
            build_quote(prs, s, n, dark_hex, slide_accent)
        else:
            build_bullets(prs, s, n, dark_hex, slide_accent)

    return prs


# ─────────────────────────────────────────────────────────────────────────────
# Social Decode Slides (appended after website slides)
# ─────────────────────────────────────────────────────────────────────────────

def parse_social_decode(decode_md: str, social_report_md: str = "") -> dict:
    """Extract key sections from social-decode.md into structured data."""
    import re

    def extract_section(md, heading):
        pattern = rf"## {re.escape(heading)}\n(.*?)(?=\n## |\Z)"
        m = re.search(pattern, md, re.DOTALL)
        return m.group(1).strip() if m else ""

    def clean_bold(text):
        """Strip **bold** markdown markers from text."""
        return re.sub(r'\*\*([^*]+)\*\*', r'\1', text)

    def bullet_lines(text, max_items=6):
        lines = []
        for l in text.split("\n"):
            l = l.strip()
            if not l or l.startswith("#"):
                continue
            # Strip leading bullet chars (including any leading **)
            l = l.lstrip("-•*").strip()
            # Strip bold markers (matched pairs first, then orphan **)
            l = clean_bold(l)
            l = l.replace('**', '').strip()
            # Skip sub-heading lines that end with colon and nothing else
            # e.g. "Dominant tone:", "Repeating language patterns:"
            if re.match(r'^[\w\s,/()-]+:\s*$', l):
                continue
            # Skip table separator rows
            if re.match(r'^[-|: ]+$', l):
                continue
            if l:
                lines.append(l)
        return lines[:max_items]

    voice_text    = extract_section(decode_md, "1. Brand Voice on LinkedIn")
    strategy_text = extract_section(decode_md, "2. Content Strategy Breakdown")
    messaging     = extract_section(decode_md, "3. Messaging Patterns")
    top3_text     = extract_section(decode_md, "7. Top 3 Posts Analysis")
    steal_text    = extract_section(decode_md, "8. What to Steal / What to Avoid")
    brief_text    = extract_section(decode_md, "9. One-Paragraph Social Positioning Brief")

    # Parse content strategy table — flexible column handling
    strategy_bullets = []
    for line in strategy_text.split("\n"):
        if not line.strip().startswith("|"):
            continue
        parts = [p.strip() for p in line.split("|") if p.strip()]
        if len(parts) < 2:
            continue
        category = parts[0]
        # Skip header and separator rows
        if re.match(r'^[-:]+$', category) or category.lower() in ("type", "category", "post type"):
            continue
        count_val = None
        pct_val = None
        for p in parts[1:]:
            if re.match(r'^\d+$', p) and count_val is None:
                count_val = p
            pct_m = re.search(r'[\d.]+\s*%', p)
            if pct_m and pct_val is None:
                pct_val = pct_m.group()
        if count_val and count_val != "0":
            label = f"{category}: {pct_val or ''} ({count_val} posts)".strip()
            strategy_bullets.append(clean_bold(label))
    if not strategy_bullets:
        strategy_bullets = bullet_lines(strategy_text)

    # Parse steal/avoid
    steal_bullets = []
    for line in steal_text.split("\n"):
        l = line.strip()
        if l.startswith("**✓") or l.startswith("**✗") or l.startswith("✓") or l.startswith("✗"):
            clean = clean_bold(re.sub(r"\*\*", "", l)).strip()
            steal_bullets.append(clean)

    # Build post URL map from social-report.md URNs
    post_urls = {}
    if social_report_md:
        for m in re.finditer(r"### Post (\d+)[^\n]*\n\*\*URN/ID:\*\*\s*`(\d+)`", social_report_md):
            num = int(m.group(1))
            urn = m.group(2)
            post_urls[num] = f"https://www.linkedin.com/feed/update/urn:li:share:{urn}/"

    # Top 3 post summaries — preserve post number for hyperlinking
    # Handles multiple formats:
    #   **Post N — Title (N likes)**\n> "quote"
    #   ### 🏆 Post N — Title (N likes)\n\n> "quote"
    top3_bullets = []
    for m in re.finditer(
        r"(?:\*\*|###\s*[^\w]*)\s*Post (\d+)[^*\n]*?(?:\*\*)?\s*[—–-]?\s*([^\n(]*?)(?:\s*\(\d+[^)]*\))?\s*\n(.*?)(?=\n(?:\*\*|###\s*[^\w]*)\s*Post |\Z)",
        top3_text, re.DOTALL
    ):
        post_num = int(m.group(1))
        title_hint = re.sub(r'^>\s*[\*_]*', '', m.group(2).strip()).strip('" *_')
        block_lines = [ll.strip() for ll in m.group(3).strip().split("\n") if ll.strip()]

        # Prefer blockquote from block content (starts with >) over title hint
        blockquote = None
        for line in block_lines:
            if line.startswith(">"):
                blockquote = re.sub(r'^>\s*[\*_]*"?', '', line).rstrip('"*_ ')
                blockquote = clean_bold(blockquote)[:120]
                break

        if blockquote and len(blockquote) > 8:
            summary = blockquote
        elif title_hint and len(title_hint) > 8:
            summary = clean_bold(title_hint)[:120]
        elif block_lines:
            first = clean_bold(block_lines[0].replace("Strategic importance:", "").replace("**", "").strip())
            summary = first[:120]
        else:
            continue

        if summary:
            top3_bullets.append({"post_num": post_num, "text": f"Post {post_num} — {summary}"})

    # Fallback if regex didn't match
    if not top3_bullets:
        for block in re.split(r"(?:\*\*|###\s*[^\w]*)\s*Post \d", top3_text):
            # Find first blockquote line > "..."
            quote_m = re.search(r'>\s*[\*_]*"?([^"\n]+)"?', block)
            if quote_m:
                first_line = clean_bold(quote_m.group(1).strip())
            else:
                first_line = clean_bold(block.strip().split("\n")[0].strip(" —*>\"#"))
            if first_line and len(first_line) > 10:
                top3_bullets.append({"post_num": None, "text": first_line[:120]})
    top3_bullets = top3_bullets[:3]

    # Messaging keywords
    msg_bullets = bullet_lines(messaging)

    return {
        "voice":     bullet_lines(voice_text, 5),
        "strategy":  strategy_bullets[:6],
        "messaging": msg_bullets,
        "top3":      top3_bullets,
        "steal":     steal_bullets[:5],
        "brief":     brief_text[:500],
        "post_urls": post_urls,
    }


def build_social_divider(prs, dark_hex, accent_hex, slide_num):
    """Section break slide: 'Social Intelligence Decode'."""
    slide = blank_slide(prs)
    fill_bg(slide, dark_hex)
    W, H = SLIDE_W, SLIDE_H

    # Accent bar
    add_rect(slide, 0, H * 0.42, W, Inches(0.06), accent_hex)

    add_text(slide, "SOCIAL INTELLIGENCE", Inches(0.5), H * 0.22, W - Inches(1),
             Inches(0.5), size=Pt(13), bold=False, color=hex_to_rgb(accent_hex))
    add_text(slide, "LinkedIn Decode", Inches(0.5), H * 0.30, W - Inches(1),
             Inches(0.9), size=Pt(44), bold=True, color=WHITE)

    # Slide number
    add_text(slide, str(slide_num), W - Inches(0.8), H - Inches(0.4),
             Inches(0.6), Inches(0.3), size=Pt(9), color=lighten(accent_hex, 0.5),
             align=PP_ALIGN.RIGHT)


def build_social_profile(prs, profile: dict, dark_hex, accent_hex, slide_num):
    """Company profile overview slide."""
    slide = blank_slide(prs)
    fill_bg(slide, "#F4F5F7")
    W, H = SLIDE_W, SLIDE_H
    panel_w = Inches(3.2)

    # Left dark panel
    add_rect(slide, 0, 0, panel_w, H, dark_hex)
    add_text(slide, "COMPANY", Inches(0.3), Inches(0.35), panel_w - Inches(0.4),
             Inches(0.3), size=Pt(9), color=hex_to_rgb(accent_hex), bold=True)
    add_text(slide, "PROFILE", Inches(0.3), Inches(0.6), panel_w - Inches(0.4),
             Inches(0.4), size=Pt(22), bold=True, color=WHITE)

    # Stats stacked in left panel
    stats = [
        ("Followers",   profile.get("follower_count", "—")),
        ("Size",        profile.get("company_size", "—")),
        ("Founded",     profile.get("founded", "—")),
        ("Funding",     profile.get("funding", "—")),
        ("Round",       profile.get("funding_round", "—")),
        ("HQ",          profile.get("headquarters", "—")),
    ]
    y = Inches(1.3)
    for label, val in stats:
        add_text(slide, label.upper(), Inches(0.3), y, panel_w - Inches(0.4),
                 Inches(0.22), size=Pt(7), color=hex_to_rgb(accent_hex), bold=True)
        add_text(slide, str(val), Inches(0.3), y + Inches(0.2), panel_w - Inches(0.4),
                 Inches(0.28), size=Pt(11), color=WHITE)
        y += Inches(0.58)

    # Right: name + tagline + description + specialties
    rx = panel_w + Inches(0.35)
    rw = W - rx - Inches(0.25)
    add_text(slide, profile.get("name", ""), rx, Inches(0.25), rw, Inches(0.45),
             size=Pt(20), bold=True, color=hex_to_rgb(dark_hex))

    tagline = profile.get("tagline", "")
    if tagline:
        # Strip emoji prefix if present
        import re as _re
        tagline_clean = _re.sub(r'^[\U00010000-\U0010ffff⚡🔥✅🚀💡🎯]+\s*', '', tagline).strip()
        add_text(slide, tagline_clean[:120], rx, Inches(0.7), rw, Inches(0.35),
                 size=Pt(9), color=hex_to_rgb(accent_hex), wrap=True)
        desc_top = Inches(1.08)
    else:
        desc_top = Inches(0.85)

    desc = profile.get("description", "")[:600]
    add_text(slide, desc, rx, desc_top, rw, Inches(2.0),
             size=Pt(9), color=MID_GRAY, wrap=True)

    specs = profile.get("specialties", "")
    if specs:
        add_text(slide, "SPECIALTIES", rx, Inches(3.2), rw, Inches(0.25),
                 size=Pt(8), bold=True, color=hex_to_rgb(accent_hex))
        add_text(slide, str(specs)[:250], rx, Inches(3.45), rw, Inches(1.6),
                 size=Pt(9), color=MID_GRAY, wrap=True)

    slide_footer(slide, slide_num, profile.get("name", "linkedin"), accent_hex)


def build_social_bullets(prs, title, subtitle, bullets, dark_hex, accent_hex, slide_num):
    """Reuse the bullets layout for social slides."""
    slide_data = {
        "title": title,
        "subtitle": subtitle,
        "body": bullets,
    }
    build_bullets(prs, slide_data, slide_num, dark_hex, accent_hex)


def build_social_quote(prs, quote, attribution, dark_hex, accent_hex, slide_num):
    slide_data = {
        "quote": quote,
        "attribution": attribution,
        "bg_color": dark_hex,
        "accent_color": accent_hex,
    }
    build_quote(prs, slide_data, slide_num, dark_hex, accent_hex)


def build_social_top3(prs, top3_items, post_urls, dark_hex, accent_hex, slide_num):
    """Top 3 Posts slide with hyperlinked Post N references."""
    slide = blank_slide(prs)
    HEADER_H = Inches(1.1)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, "#F4F5F7")
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, HEADER_H, dark_hex)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Pt(5), accent_hex)

    add_text(slide, "Top 3 Posts Analysis",
             Inches(0.5), Inches(0.15), Inches(9.0), Inches(0.85),
             size=Pt(24), bold=True, color=WHITE)
    add_text(slide, "Highest-signal posts and what they reveal",
             Inches(0.5), HEADER_H + Inches(0.08), Inches(9.0), Inches(0.35),
             size=Pt(14), color=hex_to_rgb(accent_hex))

    body_top = HEADER_H + Inches(0.5)
    for i, item in enumerate(top3_items[:3]):
        top = body_top + i * Inches(1.1)
        if top + Inches(0.9) > SLIDE_H - Inches(0.35):
            break
        text = item["text"] if isinstance(item, dict) else item
        post_num = item.get("post_num") if isinstance(item, dict) else None

        # Accent dot
        add_rect(slide, Inches(0.5), top + Inches(0.18), Inches(0.07), Inches(0.26), accent_hex)

        add_text_with_links(slide, text,
                            Inches(0.7), top, Inches(9.0), Inches(1.0),
                            size=Pt(13), color=MID_GRAY, post_urls=post_urls)

    add_rect(slide, Inches(0), SLIDE_H - Pt(5), SLIDE_W, Pt(5), accent_hex)
    add_text(slide, str(slide_num), SLIDE_W - Inches(0.55), SLIDE_H - Inches(0.38),
             Inches(0.4), Inches(0.35), size=Pt(9), color=RGBColor(0xAA, 0xAA, 0xBB),
             align=PP_ALIGN.RIGHT)


def build_social_slides(prs, social_report_path, social_decode_path,
                        dark_hex, accent_hex, start_slide_num):
    """Append social decode slides to an existing presentation."""
    import json as _json

    n = start_slide_num

    # Section divider
    build_social_divider(prs, dark_hex, accent_hex, n); n += 1

    # Parse profile from social-report.md JSON block
    profile = {}
    if social_report_path and Path(social_report_path).exists():
        report_text = Path(social_report_path).read_text(encoding="utf-8")
        import re
        m = re.search(r"## Raw Profile JSON\s*```json\s*(\{.*?\})\s*```", report_text, re.DOTALL)
        if m:
            try:
                profile = _json.loads(m.group(1))
            except Exception:
                pass

        # Also extract posts for top posts slide
        posts = []
        pm = re.search(r"## Raw Posts JSON\s*```json\s*(\[.*?\])\s*```", report_text, re.DOTALL)
        if pm:
            try:
                posts = _json.loads(pm.group(1))
            except Exception:
                pass

    # Company profile slide
    build_social_profile(prs, profile, dark_hex, accent_hex, n); n += 1

    # Also read social-report.md text for URN extraction
    report_text_for_parse = Path(social_report_path).read_text(encoding="utf-8") if (social_report_path and Path(social_report_path).exists()) else ""

    # Parse social-decode.md
    decoded = {}
    if social_decode_path and Path(social_decode_path).exists():
        decode_text = Path(social_decode_path).read_text(encoding="utf-8")
        decoded = parse_social_decode(decode_text, report_text_for_parse)

    post_urls = decoded.get("post_urls", {})

    # Brand voice + content strategy
    voice = decoded.get("voice", ["No decode available — run with ANTHROPIC_API_KEY"])
    build_social_bullets(prs, "Brand Voice on LinkedIn", "Tone, language patterns, and consistency",
                         voice, dark_hex, accent_hex, n); n += 1

    strategy = decoded.get("strategy", [])
    if strategy:
        build_social_bullets(prs, "Content Strategy Breakdown", "Post type mix across last 10 posts",
                             strategy, dark_hex, accent_hex, n); n += 1

    # Messaging patterns
    messaging = decoded.get("messaging", [])
    if messaging:
        build_social_bullets(prs, "Messaging Patterns", "Problems named, outcomes promised, keywords owned",
                             messaging, dark_hex, accent_hex, n); n += 1

    # Top 3 posts — with hyperlinks on "Post N" references
    top3 = decoded.get("top3", [])
    if top3:
        build_social_top3(prs, top3, post_urls, dark_hex, accent_hex, n); n += 1

    # Steal / Avoid
    steal = decoded.get("steal", [])
    if steal:
        build_social_bullets(prs, "What to Steal / What to Avoid",
                             "Tactics worth replicating and patterns to skip",
                             steal, dark_hex, accent_hex, n); n += 1

    # Positioning brief quote
    brief = decoded.get("brief", "")
    if brief:
        build_social_quote(prs, brief, profile.get("name", "linkedin"),
                           dark_hex, accent_hex, n); n += 1

    return prs


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--domain",      required=True)
    parser.add_argument("--run-id",      required=True)
    parser.add_argument("--repo1",       required=False, default="")
    parser.add_argument("--local-path",  required=False, default="",
                        help="Local path to website-decoder repo (skips GitHub fetch)")
    parser.add_argument("--social-path", required=False, default="",
                        help="Local path to social-decoder output dir (adds social slides)")
    args = parser.parse_args()

    local_repo = Path(args.local_path) if args.local_path else None

    if local_repo:
        json_path = local_repo / "output" / "decodes" / args.domain / "ppt_content.json"
        print(f"Loading local: {json_path}")
        if not json_path.exists():
            print(f"ERROR: {json_path} not found")
            sys.exit(1)
        content = json.loads(json_path.read_text(encoding="utf-8"))
        screenshots_base = str(local_repo / "output" / "screenshots" / args.domain / args.run_id)
        cache_dir = local_repo / "output" / "screenshots" / args.domain / args.run_id
    else:
        raw_base = f"https://raw.githubusercontent.com/{args.repo1}/main"
        json_url = f"{raw_base}/output/decodes/{args.domain}/ppt_content.json"
        print(f"Fetching: {json_url}")
        data = fetch_bytes(json_url)
        if data is None:
            print("ERROR: ppt_content.json not found")
            sys.exit(1)
        content = json.loads(data.decode("utf-8"))
        screenshots_base = f"{raw_base}/output/screenshots/{args.domain}/{args.run_id}"
        cache_dir = Path("tmp") / args.domain / args.run_id

    print(f"Building website slides: {args.domain} / {args.run_id}")
    prs = build_presentation(content, cache_dir, screenshots_base)

    # Append social slides if --social-path provided
    if args.social_path:
        social_dir = Path(args.social_path)
        report_path = social_dir / "social-report.md"
        decode_path = social_dir / "social-decode.md"

        # Extract brand colors from cover slide
        slides = content.get("slides", [])
        cover  = next((s for s in slides if s.get("id") == 1), slides[0] if slides else {})
        dark_hex   = cover.get("bg_color",    "#0C141E")
        accent_hex = cover.get("accent_color", "#53BDE4")

        start_n = len(prs.slides) + 1
        print(f"Appending social slides from: {social_dir}")
        build_social_slides(prs, str(report_path), str(decode_path),
                            dark_hex, accent_hex, start_n)

    out_dir = Path("output")
    out_dir.mkdir(exist_ok=True)
    suffix = "-combined" if args.social_path else ""
    out_path = out_dir / f"{args.domain}-{args.run_id}{suffix}.pptx"
    prs.save(str(out_path))
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    main()
